from __future__ import annotations

import io
import json
import logging
import random
import re
import shutil
import subprocess
import tempfile
import time
import textwrap
import uuid
from dataclasses import dataclass
from datetime import date, datetime, timedelta, timezone
from pathlib import Path
from typing import Iterable
from urllib.parse import urlparse

from google.auth import default as google_auth_default
from google.api_core import exceptions as google_api_exceptions
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload
from google.cloud import bigquery
from google.cloud import storage
import vertexai
from vertexai.generative_models import GenerativeModel, Part
from vertexai.language_models import TextEmbeddingInput, TextEmbeddingModel
from pypdf import PdfReader
from pptx import Presentation

from .config import (
    BQ_DATASET,
    BQ_TABLE,
    CHUNK_OVERLAP,
    CHUNK_SIZE,
    EMBED_BATCH_SIZE,
    EMBED_MODEL,
    GCP_LOCATION,
    GCP_PROJECT_ID,
    GCS_BUCKET,
    GEN_FALLBACK_MODELS,
    GEN_MODEL,
    MEDIA_FALLBACK_MODELS,
    MEDIA_GEN_MODEL,
    VIDEO_ENABLE_VISUAL_ANALYSIS,
    VERTEX_MAX_RETRIES,
    VERTEX_RETRY_INITIAL_SECONDS,
    VERTEX_RETRY_MAX_SECONDS,
)


@dataclass
class SlideChunk:
    source_id: str
    source_uri: str
    source_name: str
    source_system: str
    title: str | None
    slide_number: int
    chunk_index: int
    content_type: str
    modalities: list[str]
    detected_date: str | None
    media_start_seconds: float | None
    media_end_seconds: float | None
    speech_style: str | None
    word_timestamps_json: str | None
    chunk_text: str


class SlideRAGService:
    def __init__(self) -> None:
        self.bq = bigquery.Client(project=GCP_PROJECT_ID)
        self.storage = storage.Client(project=GCP_PROJECT_ID)
        vertexai.init(project=GCP_PROJECT_ID, location=GCP_LOCATION)
        self.embedder = TextEmbeddingModel.from_pretrained(EMBED_MODEL)
        self.generator = GenerativeModel(GEN_MODEL)
        self.media_models = self._build_model_chain(MEDIA_GEN_MODEL, MEDIA_FALLBACK_MODELS)
        text_fallbacks = list(GEN_FALLBACK_MODELS)
        if MEDIA_GEN_MODEL not in text_fallbacks and MEDIA_GEN_MODEL != GEN_MODEL:
            text_fallbacks.append(MEDIA_GEN_MODEL)
        for model_name in MEDIA_FALLBACK_MODELS:
            if model_name not in text_fallbacks and model_name != GEN_MODEL:
                text_fallbacks.append(model_name)
        self.text_models = self._build_model_chain(GEN_MODEL, text_fallbacks)
        self.logger = logging.getLogger(__name__)

    def upload_to_gcs(self, filename: str, content: bytes) -> str:
        safe_name = f"{uuid.uuid4()}-{Path(filename).name}"
        bucket = self.storage.bucket(GCS_BUCKET)
        blob = bucket.blob(f"slides/{safe_name}")
        blob.upload_from_string(content)
        return f"gs://{GCS_BUCKET}/slides/{safe_name}"

    def ingest_bytes(self, filename: str, content: bytes, source_system: str = "upload") -> dict:
        self.logger.info(
            "Ingest start filename=%s source_system=%s size_bytes=%d",
            filename,
            source_system,
            len(content),
        )
        source_uri = self.upload_to_gcs(filename, content)
        self.logger.info("Uploaded to GCS source_uri=%s", source_uri)
        chunks = self.extract_chunks(filename, content, source_uri, source_system)
        self.logger.info("Chunk extraction complete source_uri=%s chunks=%d", source_uri, len(chunks))
        inserted = self.upsert_chunks(chunks)
        self.logger.info("BigQuery upsert complete source_uri=%s inserted=%d", source_uri, inserted)
        return {
            "source_id": self._source_id_from_uri(source_uri),
            "source_uri": source_uri,
            "chunks_inserted": inserted,
        }

    def extract_chunks(
        self, filename: str, content: bytes, source_uri: str, source_system: str = "upload"
    ) -> list[SlideChunk]:
        suffix = Path(filename).suffix.lower()
        if suffix == ".pptx":
            return self._extract_from_pptx(filename, content, source_uri, source_system)
        if suffix == ".pdf":
            return self._extract_from_pdf(filename, content, source_uri, source_system)
        if suffix in {".png", ".jpg", ".jpeg", ".webp"}:
            return self._extract_from_image(filename, content, source_uri, source_system, 1)
        if suffix in {".mp3", ".wav", ".m4a"}:
            return self._extract_from_audio(filename, content, source_uri, source_system)
        if suffix in {".mp4", ".mov", ".webm", ".mkv"}:
            return self._extract_from_video(filename, content, source_uri, source_system)
        raise ValueError(
            "Supported: .pptx, .pdf, .png, .jpg, .jpeg, .webp, .mp3, .wav, .m4a, .mp4, .mov, .webm, .mkv."
        )

    def upsert_chunks(self, chunks: list[SlideChunk]) -> int:
        if not chunks:
            return 0
        embeddings = self._embed_documents([c.chunk_text for c in chunks])
        rows = []
        for chunk, embedding in zip(chunks, embeddings, strict=False):
            rows.append(
                {
                    "chunk_id": str(uuid.uuid4()),
                    "source_id": chunk.source_id,
                    "source_uri": chunk.source_uri,
                    "source_name": chunk.source_name,
                    "source_system": chunk.source_system,
                    "title": chunk.title,
                    "slide_number": chunk.slide_number,
                    "chunk_index": chunk.chunk_index,
                    "content_type": chunk.content_type,
                    "modalities": chunk.modalities,
                    "detected_date": chunk.detected_date,
                    "media_start_seconds": chunk.media_start_seconds,
                    "media_end_seconds": chunk.media_end_seconds,
                    "speech_style": chunk.speech_style,
                    "word_timestamps_json": chunk.word_timestamps_json,
                    "chunk_text": chunk.chunk_text,
                    "embedding": embedding,
                    "created_at": datetime.now(timezone.utc).isoformat(),
                }
            )

        table = f"{GCP_PROJECT_ID}.{BQ_DATASET}.{BQ_TABLE}"
        errors = self.bq.insert_rows_json(table, rows)
        if errors:
            raise RuntimeError(f"BigQuery insert errors: {errors}")
        return len(rows)

    def answer_question(self, question: str, top_k: int = 5) -> dict:
        query_embedding = self._embed_query(question)
        date_start, date_end = self._extract_date_filter(question)
        search_k = min(max(top_k * 5, top_k), 100)
        sql = f"""
        WITH hits AS (
          SELECT
            base.source_uri AS source_uri,
            base.source_id AS source_id,
            base.source_name AS source_name,
            base.source_system AS source_system,
            base.title AS title,
            base.slide_number AS slide_number,
            base.chunk_index AS chunk_index,
            base.content_type AS content_type,
            base.modalities AS modalities,
            base.detected_date AS detected_date,
            base.media_start_seconds AS media_start_seconds,
            base.media_end_seconds AS media_end_seconds,
            base.speech_style AS speech_style,
            base.chunk_text AS chunk_text,
            distance
          FROM VECTOR_SEARCH(
            TABLE `{GCP_PROJECT_ID}.{BQ_DATASET}.{BQ_TABLE}`,
            'embedding',
            (SELECT @query_embedding AS embedding),
            top_k => @search_k,
            distance_type => 'COSINE'
          )
        )
        SELECT *
        FROM hits
        WHERE (@date_start IS NULL OR detected_date >= @date_start)
          AND (@date_end IS NULL OR detected_date <= @date_end)
        ORDER BY distance
        LIMIT @top_k
        """
        params = [
            bigquery.ArrayQueryParameter("query_embedding", "FLOAT64", query_embedding),
            bigquery.ScalarQueryParameter("search_k", "INT64", search_k),
            bigquery.ScalarQueryParameter("top_k", "INT64", top_k),
            bigquery.ScalarQueryParameter("date_start", "DATE", date_start),
            bigquery.ScalarQueryParameter("date_end", "DATE", date_end),
        ]
        job_config = bigquery.QueryJobConfig(query_parameters=params)
        hits = list(self.bq.query(sql, job_config=job_config).result())

        if not hits:
            return {"answer": "No relevant content found in indexed slides.", "citations": []}

        context_lines = []
        citations = []
        for h in hits:
            if h.media_start_seconds is not None:
                end_s = h.media_end_seconds if h.media_end_seconds is not None else h.media_start_seconds
                context_lines.append(f"[source={h.source_uri}, t={h.media_start_seconds:.2f}-{end_s:.2f}s] {h.chunk_text}")
            else:
                context_lines.append(f"[source={h.source_uri}, slide={h.slide_number}] {h.chunk_text}")
            citations.append(
                {
                    "source_uri": h.source_uri,
                    "source_signed_url": self._maybe_sign_gs_uri(h.source_uri),
                    "source_id": h.source_id,
                    "source_name": h.source_name,
                    "source_system": h.source_system,
                    "slide_number": h.slide_number,
                    "chunk_index": h.chunk_index,
                    "title": h.title,
                    "content_type": h.content_type,
                    "modalities": h.modalities,
                    "detected_date": h.detected_date,
                    "media_start_seconds": h.media_start_seconds,
                    "media_end_seconds": h.media_end_seconds,
                    "speech_style": h.speech_style,
                    "distance": h.distance,
                }
            )

        prompt = textwrap.dedent(
            f"""
            You are a precise assistant answering questions from slide content.
            Use only the provided context. If context is insufficient, say so clearly.
            Question: {question}

            Context:
            {'\n'.join(context_lines)}
            """
        )
        response = self._run_with_ai_retries(
            lambda: self.generator.generate_content(prompt),
            operation="answer generation",
        )
        answer = response.text if hasattr(response, "text") and response.text else str(response)
        return {"answer": answer, "citations": citations}

    def ingest_drive(
        self,
        file_id: str | None = None,
        folder_id: str | None = None,
        max_files: int = 50,
    ) -> dict:
        if not file_id and not folder_id:
            raise ValueError("Provide file_id or folder_id.")
        if file_id and folder_id:
            raise ValueError("Provide only one of file_id or folder_id.")

        drive = self._drive_service()
        candidates = [self._drive_file_meta(drive, file_id)] if file_id else self._drive_list_folder(
            drive, folder_id, max_files
        )

        indexed = []
        skipped = []
        total_chunks = 0
        for item in candidates:
            normalized = self._normalize_drive_file(item)
            if not normalized:
                skipped.append({"file_id": item["id"], "name": item["name"], "reason": "unsupported file type"})
                continue

            target_name, content = self._download_drive_file(drive, item, normalized)
            result = self.ingest_bytes(target_name, content, source_system="drive")
            total_chunks += result["chunks_inserted"]
            indexed.append(
                {
                    "file_id": item["id"],
                    "name": item["name"],
                    "source_id": result["source_id"],
                    "source_uri": result["source_uri"],
                    "chunks_inserted": result["chunks_inserted"],
                }
            )

        return {
            "indexed_files": len(indexed),
            "indexed_chunks": total_chunks,
            "indexed": indexed,
            "skipped": skipped,
        }

    def _extract_from_pptx(
        self, filename: str, content: bytes, source_uri: str, source_system: str
    ) -> list[SlideChunk]:
        deck = Presentation(io.BytesIO(content))
        chunks: list[SlideChunk] = []
        title = None
        source_name = Path(filename).name

        for idx, slide in enumerate(deck.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text.strip())
                if hasattr(shape, "image") and getattr(shape, "image", None):
                    image_caption = self._describe_binary(
                        shape.image.blob,
                        "image/jpeg",
                        "Describe this slide image in 3-6 concise bullets.",
                    )
                    texts.append(f"[Image summary] {image_caption}")
            slide_text = "\n".join(t for t in texts if t)
            if idx == 1 and texts:
                title = texts[0]
            chunks.extend(
                self._chunk_slide_text(
                    text=slide_text,
                    source_uri=source_uri,
                    source_name=source_name,
                    source_system=source_system,
                    title=title,
                    slide_number=idx,
                    content_type="slide",
                    modalities=["text", "image"],
                )
            )
        return chunks

    def _extract_from_pdf(
        self, filename: str, content: bytes, source_uri: str, source_system: str
    ) -> list[SlideChunk]:
        reader = PdfReader(io.BytesIO(content))
        chunks: list[SlideChunk] = []
        title = None
        source_name = Path(filename).name

        for idx, page in enumerate(reader.pages, start=1):
            page_text = (page.extract_text() or "").strip()
            image_texts = []
            try:
                for img in getattr(page, "images", []):
                    image_texts.append(
                        self._describe_binary(
                            img.data,
                            "image/jpeg",
                            "Describe this page image in 2-5 concise bullets.",
                        )
                    )
            except Exception:
                pass
            if image_texts:
                page_text = f"{page_text}\n[Page image summaries]\n" + "\n".join(image_texts)
            if idx == 1 and page_text:
                title = page_text.splitlines()[0][:180]
            chunks.extend(
                self._chunk_slide_text(
                    text=page_text,
                    source_uri=source_uri,
                    source_name=source_name,
                    source_system=source_system,
                    title=title,
                    slide_number=idx,
                    content_type="pdf_page",
                    modalities=["text", "image"],
                )
            )
        return chunks

    def _extract_from_image(
        self,
        filename: str,
        content: bytes,
        source_uri: str,
        source_system: str,
        slide_number: int,
    ) -> list[SlideChunk]:
        source_name = Path(filename).name
        suffix = Path(filename).suffix.lower()
        mime = "image/jpeg"
        if suffix == ".png":
            mime = "image/png"
        if suffix == ".webp":
            mime = "image/webp"
        text = self._describe_binary(
            content,
            mime,
            "Describe the image in detail. Include visible text, entities, charts, and key scene context.",
        )
        return self._chunk_slide_text(
            text=text,
            source_uri=source_uri,
            source_name=source_name,
            source_system=source_system,
            title=source_name,
            slide_number=slide_number,
            content_type="image",
            modalities=["image", "text"],
        )

    def _extract_from_audio(
        self, filename: str, content: bytes, source_uri: str, source_system: str
    ) -> list[SlideChunk]:
        source_name = Path(filename).name
        suffix = Path(filename).suffix.lower()
        mime_map = {
            ".mp3": "audio/mpeg",
            ".wav": "audio/wav",
            ".m4a": "audio/mp4",
        }
        mime = mime_map.get(suffix, "audio/mpeg")
        transcript = self._transcribe_media_with_timestamps(
            content,
            mime,
            include_visual_analysis=False,
            source_uri=source_uri,
        )
        return self._media_segments_to_chunks(
            transcript=transcript,
            source_uri=source_uri,
            source_name=source_name,
            source_system=source_system,
            content_type="audio",
            modalities=["audio", "text"],
        )

    def _extract_from_video(
        self, filename: str, content: bytes, source_uri: str, source_system: str
    ) -> list[SlideChunk]:
        source_name = Path(filename).name
        suffix = Path(filename).suffix.lower()
        mime_map = {
            ".mp4": "video/mp4",
            ".mov": "video/quicktime",
            ".webm": "video/webm",
            ".mkv": "video/x-matroska",
        }
        mime = mime_map.get(suffix, "video/mp4")
        transcript = self._transcribe_media_with_timestamps(
            content,
            mime,
            include_visual_analysis=VIDEO_ENABLE_VISUAL_ANALYSIS,
            source_uri=source_uri,
        )
        return self._media_segments_to_chunks(
            transcript=transcript,
            source_uri=source_uri,
            source_name=source_name,
            source_system=source_system,
            content_type="video",
            modalities=["video", "audio", "text"],
        )

    def _media_segments_to_chunks(
        self,
        transcript: dict,
        source_uri: str,
        source_name: str,
        source_system: str,
        content_type: str,
        modalities: list[str],
    ) -> list[SlideChunk]:
        source_id = self._source_id_from_uri(source_uri)
        chunks: list[SlideChunk] = []

        summary = str(transcript.get("summary", "")).strip()
        if summary:
            chunks.append(
                SlideChunk(
                    source_id=source_id,
                    source_uri=source_uri,
                    source_name=source_name,
                    source_system=source_system,
                    title=source_name,
                    slide_number=1,
                    chunk_index=1,
                    content_type=content_type,
                    modalities=modalities,
                    detected_date=self._extract_date_from_text(summary) or self._extract_date_from_text(source_name),
                    media_start_seconds=None,
                    media_end_seconds=None,
                    speech_style=None,
                    word_timestamps_json=None,
                    chunk_text=f"[Summary] {summary}",
                )
            )

        start_index = len(chunks) + 1
        segments = self._expand_sparse_media_segments(transcript)
        for i, segment in enumerate(segments, start=start_index):
            segment_text = str(segment.get("text", "")).strip()
            if not segment_text:
                continue
            start_s = self._safe_float(segment.get("start_seconds"))
            end_s = self._safe_float(segment.get("end_seconds"))
            style = str(segment.get("style", "")).strip() or None

            words = []
            for word_item in segment.get("words", []) or []:
                if not isinstance(word_item, dict):
                    continue
                word = str(word_item.get("word", "")).strip()
                if not word:
                    continue
                words.append(
                    {
                        "word": word,
                        "start_seconds": self._safe_float(word_item.get("start_seconds")),
                        "end_seconds": self._safe_float(word_item.get("end_seconds")),
                        "style": str(word_item.get("style", "")).strip() or style,
                    }
                )
            word_timestamps_json = json.dumps(words, separators=(",", ":")) if words else None

            chunk_text = segment_text
            if start_s is not None:
                end_val = end_s if end_s is not None else start_s
                chunk_text = f"[t={start_s:.2f}-{end_val:.2f}s] {segment_text}"
            if style:
                chunk_text = f"{chunk_text} [style={style}]"

            chunks.append(
                SlideChunk(
                    source_id=source_id,
                    source_uri=source_uri,
                    source_name=source_name,
                    source_system=source_system,
                    title=source_name,
                    slide_number=1,
                    chunk_index=i,
                    content_type=content_type,
                    modalities=modalities,
                    detected_date=self._extract_date_from_text(segment_text) or self._extract_date_from_text(source_name),
                    media_start_seconds=start_s,
                    media_end_seconds=end_s,
                    speech_style=style,
                    word_timestamps_json=word_timestamps_json,
                    chunk_text=chunk_text,
                )
            )

        if chunks:
            return chunks

        fallback_text = str(transcript.get("text", "")).strip()
        if not fallback_text:
            return []
        return self._chunk_slide_text(
            text=f"[Transcript] {fallback_text}",
            source_uri=source_uri,
            source_name=source_name,
            source_system=source_system,
            title=source_name,
            slide_number=1,
            content_type=content_type,
            modalities=modalities,
        )

    def _expand_sparse_media_segments(self, transcript: dict) -> list[dict]:
        raw_segments = transcript.get("segments") or []
        if len(raw_segments) > 1:
            return raw_segments

        base = raw_segments[0] if raw_segments else {}
        text = str(base.get("text") or transcript.get("text") or "").strip()
        if not text:
            return raw_segments

        start_seconds = self._safe_float(base.get("start_seconds")) or 0.0
        end_seconds = self._safe_float(base.get("end_seconds"))
        style = str(base.get("style", "")).strip() or None
        words = text.split()

        # If timing already looks meaningful, keep as-is.
        if end_seconds is not None and end_seconds > start_seconds and len(words) < 90:
            return raw_segments
        if len(words) < 35:
            return raw_segments

        sentence_units = [s.strip() for s in re.split(r"(?<=[.!?])\s+", text) if s.strip()]
        units = sentence_units if len(sentence_units) >= 3 else []
        if not units:
            chunk_size_words = 28
            units = [" ".join(words[i : i + chunk_size_words]) for i in range(0, len(words), chunk_size_words)]

        segments: list[dict] = []
        cursor = start_seconds
        total_words = max(1, sum(len(u.split()) for u in units))
        target_total = end_seconds - start_seconds if end_seconds and end_seconds > start_seconds else None

        for idx, unit in enumerate(units, start=1):
            unit_words = len(unit.split())
            if target_total is not None:
                duration = max(1.5, target_total * (unit_words / total_words))
            else:
                # Approximate speech pace when no end timestamp is available.
                duration = max(1.5, unit_words / 2.6)
            seg_start = cursor
            seg_end = seg_start + duration
            segments.append(
                {
                    "start_seconds": seg_start,
                    "end_seconds": seg_end,
                    "text": unit,
                    "style": style,
                    "words": [],
                    "timing_estimated": True,
                    "segment_index": idx,
                }
            )
            cursor = seg_end
        return segments

    def _chunk_slide_text(
        self,
        text: str,
        source_uri: str,
        source_name: str,
        source_system: str,
        title: str | None,
        slide_number: int,
        content_type: str,
        modalities: list[str],
    ) -> list[SlideChunk]:
        text = " ".join(text.split())
        if not text:
            return []

        source_id = self._source_id_from_uri(source_uri)
        chunks = []
        for chunk_index, segment in enumerate(self._windowed_chunks(text), start=1):
            chunks.append(
                SlideChunk(
                    source_id=source_id,
                    source_uri=source_uri,
                    source_name=source_name,
                    source_system=source_system,
                    title=title,
                    slide_number=slide_number,
                    chunk_index=chunk_index,
                    content_type=content_type,
                    modalities=modalities,
                    detected_date=self._extract_date_from_text(segment) or self._extract_date_from_text(source_name),
                    media_start_seconds=None,
                    media_end_seconds=None,
                    speech_style=None,
                    word_timestamps_json=None,
                    chunk_text=segment,
                )
            )
        return chunks

    def _windowed_chunks(self, text: str) -> Iterable[str]:
        step = max(1, CHUNK_SIZE - CHUNK_OVERLAP)
        for start in range(0, len(text), step):
            chunk = text[start : start + CHUNK_SIZE].strip()
            if chunk:
                yield chunk

    def list_documents(self, limit: int = 100) -> list[dict]:
        sql = f"""
        SELECT
          COALESCE(source_id, source_uri) AS source_id,
          ANY_VALUE(source_uri) AS source_uri,
          ANY_VALUE(source_name) AS source_name,
          ANY_VALUE(source_system) AS source_system,
          ANY_VALUE(content_type) AS content_type,
          ARRAY(
            SELECT DISTINCT m
            FROM UNNEST(ARRAY_CONCAT_AGG(IFNULL(modalities, []))) AS m
          ) AS modalities,
          MIN(detected_date) AS first_detected_date,
          MAX(detected_date) AS last_detected_date,
          COUNT(*) AS chunks
        FROM `{GCP_PROJECT_ID}.{BQ_DATASET}.{BQ_TABLE}`
        GROUP BY COALESCE(source_id, source_uri)
        ORDER BY MAX(created_at) DESC
        LIMIT @limit
        """
        params = [bigquery.ScalarQueryParameter("limit", "INT64", limit)]
        rows = self.bq.query(sql, job_config=bigquery.QueryJobConfig(query_parameters=params)).result()
        return [dict(r) for r in rows]

    def get_document(self, source_id: str) -> dict:
        sql = f"""
        SELECT
          source_id, source_uri, source_name, source_system, title,
          slide_number, chunk_index, content_type, modalities, detected_date, chunk_text,
          media_start_seconds, media_end_seconds, speech_style, word_timestamps_json
        FROM `{GCP_PROJECT_ID}.{BQ_DATASET}.{BQ_TABLE}`
        WHERE source_id = @source_id
        ORDER BY slide_number, chunk_index
        """
        params = [bigquery.ScalarQueryParameter("source_id", "STRING", source_id)]
        rows = list(self.bq.query(sql, job_config=bigquery.QueryJobConfig(query_parameters=params)).result())
        if not rows:
            fallback_sql = f"""
            SELECT
              source_id, source_uri, source_name, source_system, title,
              slide_number, chunk_index, content_type, modalities, detected_date, chunk_text,
              media_start_seconds, media_end_seconds, speech_style, word_timestamps_json
            FROM `{GCP_PROJECT_ID}.{BQ_DATASET}.{BQ_TABLE}`
            WHERE source_uri = @source_uri
            ORDER BY slide_number, chunk_index
            """
            fallback_params = [bigquery.ScalarQueryParameter("source_uri", "STRING", source_id)]
            rows = list(
                self.bq.query(
                    fallback_sql,
                    job_config=bigquery.QueryJobConfig(query_parameters=fallback_params),
                ).result()
            )
        if not rows and re.fullmatch(r"[0-9a-fA-F-]{36}", source_id):
            upload_id_sql = f"""
            SELECT
              source_id, source_uri, source_name, source_system, title,
              slide_number, chunk_index, content_type, modalities, detected_date, chunk_text,
              media_start_seconds, media_end_seconds, speech_style, word_timestamps_json
            FROM `{GCP_PROJECT_ID}.{BQ_DATASET}.{BQ_TABLE}`
            WHERE REGEXP_CONTAINS(source_uri, @pattern)
            ORDER BY slide_number, chunk_index
            """
            pattern = rf"/slides/{re.escape(source_id)}-"
            rows = list(
                self.bq.query(
                    upload_id_sql,
                    job_config=bigquery.QueryJobConfig(
                        query_parameters=[bigquery.ScalarQueryParameter("pattern", "STRING", pattern)]
                    ),
                ).result()
            )
        if not rows:
            raise ValueError(f"No document found for source_id={source_id}")

        sections = []
        current_slide = None
        current_texts: list[str] = []
        for r in rows:
            if current_slide is None:
                current_slide = r.slide_number
            if r.slide_number != current_slide:
                sections.append({"slide_number": current_slide, "text": " ".join(current_texts).strip()})
                current_slide = r.slide_number
                current_texts = []
            current_texts.append(r.chunk_text)
        if current_slide is not None:
            sections.append({"slide_number": current_slide, "text": " ".join(current_texts).strip()})

        full_text = "\n\n".join(
            f"[slide {s['slide_number']}]\n{s['text']}" if s["slide_number"] else s["text"] for s in sections
        ).strip()
        head = rows[0]
        return {
            "source_id": head.source_id or self._source_id_from_uri(head.source_uri),
            "source_uri": head.source_uri,
            "source_signed_url": self._maybe_sign_gs_uri(head.source_uri),
            "source_name": head.source_name,
            "source_system": head.source_system,
            "title": head.title,
            "content_type": head.content_type,
            "sections": sections,
            "full_text": full_text,
            "chunk_count": len(rows),
        }

    def find_word_occurrences(self, source_id: str, word: str, max_hits: int = 50) -> dict:
        normalized_target = re.sub(r"[^a-z0-9']+", "", word.lower())
        if not normalized_target:
            raise ValueError("Word must contain letters or numbers.")

        sql = f"""
        SELECT
          source_id, source_uri, source_name, content_type, chunk_index,
          chunk_text, media_start_seconds, media_end_seconds, speech_style, word_timestamps_json
        FROM `{GCP_PROJECT_ID}.{BQ_DATASET}.{BQ_TABLE}`
        WHERE source_id = @source_id
          AND content_type IN ('audio', 'video')
        ORDER BY chunk_index
        """
        rows = list(
            self.bq.query(
                sql,
                job_config=bigquery.QueryJobConfig(
                    query_parameters=[bigquery.ScalarQueryParameter("source_id", "STRING", source_id)]
                ),
            ).result()
        )
        if not rows:
            raise ValueError(f"No indexed audio/video content found for source_id={source_id}")
        source_signed_url = self._maybe_sign_gs_uri(rows[0].source_uri)

        occurrences: list[dict] = []
        for row in rows:
            row_style = row.speech_style
            words_payload = row.word_timestamps_json
            items = []
            if words_payload:
                try:
                    items = json.loads(words_payload)
                except json.JSONDecodeError:
                    items = []
            for item in items:
                if not isinstance(item, dict):
                    continue
                candidate = re.sub(r"[^a-z0-9']+", "", str(item.get("word", "")).lower())
                if candidate != normalized_target:
                    continue
                item_start = self._safe_float(item.get("start_seconds"))
                item_end = self._safe_float(item.get("end_seconds"))
                occurrences.append(
                    {
                        "word": str(item.get("word", "")).strip() or word,
                        "start_seconds": item_start if item_start is not None else row.media_start_seconds,
                        "end_seconds": item_end if item_end is not None else row.media_end_seconds,
                        "timecode": self._format_timecode(item_start if item_start is not None else row.media_start_seconds),
                        "style": str(item.get("style", "")).strip() or row_style,
                        "chunk_index": row.chunk_index,
                        "chunk_text": row.chunk_text,
                        "content_type": row.content_type,
                        "playback_url": self._build_playback_url(
                            source_signed_url,
                            item_start if item_start is not None else row.media_start_seconds,
                        ),
                    }
                )
                if len(occurrences) >= max_hits:
                    break
            if not items:
                chunk_tokens = [re.sub(r"[^a-z0-9']+", "", t.lower()) for t in row.chunk_text.split()]
                if normalized_target in chunk_tokens:
                    occurrences.append(
                        {
                            "word": word,
                            "start_seconds": row.media_start_seconds,
                            "end_seconds": row.media_end_seconds,
                            "timecode": self._format_timecode(row.media_start_seconds),
                            "style": row_style,
                            "chunk_index": row.chunk_index,
                            "chunk_text": row.chunk_text,
                            "content_type": row.content_type,
                            "playback_url": self._build_playback_url(source_signed_url, row.media_start_seconds),
                        }
                    )
                    if len(occurrences) >= max_hits:
                        break
            if len(occurrences) >= max_hits:
                break

        return {
            "source_id": rows[0].source_id,
            "source_uri": rows[0].source_uri,
            "source_signed_url": source_signed_url,
            "source_name": rows[0].source_name,
            "search_word": word,
            "occurrence_count": len(occurrences),
            "occurrences": occurrences,
        }

    def get_media_source(self, source_id: str) -> dict:
        doc = self.get_document(source_id)
        signed_url = self._maybe_sign_gs_uri(doc["source_uri"])
        return {
            "source_id": doc["source_id"],
            "source_name": doc["source_name"],
            "source_uri": doc["source_uri"],
            "source_signed_url": signed_url,
            "content_type": doc["content_type"],
        }

    def list_document_images(self, source_id: str, max_images: int = 50) -> dict:
        doc = self.get_document(source_id)
        source_uri = doc["source_uri"]
        source_name = doc["source_name"] or "source"
        suffix = Path(source_name).suffix.lower()

        if suffix in {".png", ".jpg", ".jpeg", ".webp"}:
            return {
                "source_id": doc["source_id"],
                "images": [self._signed_asset(source_uri, "source-image")],
            }

        content = self._download_gcs_uri(source_uri)
        source_id_norm = doc["source_id"]
        images: list[dict] = []

        if suffix == ".pptx":
            deck = Presentation(io.BytesIO(content))
            for slide_number, slide in enumerate(deck.slides, start=1):
                image_index = 0
                for shape in slide.shapes:
                    if not (hasattr(shape, "image") and getattr(shape, "image", None)):
                        continue
                    image_index += 1
                    ext = (shape.image.ext or "jpg").lower()
                    if ext == "jpeg":
                        ext = "jpg"
                    blob_path = f"extracted/{source_id_norm}/slide-{slide_number}-img-{image_index}.{ext}"
                    gs_uri = self._upload_extracted(blob_path, shape.image.blob)
                    images.append(self._signed_asset(gs_uri, f"slide-{slide_number}-img-{image_index}"))
                    if len(images) >= max_images:
                        break
                if len(images) >= max_images:
                    break

        elif suffix == ".pdf":
            reader = PdfReader(io.BytesIO(content))
            for page_number, page in enumerate(reader.pages, start=1):
                page_images = getattr(page, "images", [])
                for image_index, img in enumerate(page_images, start=1):
                    blob_path = f"extracted/{source_id_norm}/page-{page_number}-img-{image_index}.jpg"
                    gs_uri = self._upload_extracted(blob_path, img.data)
                    images.append(self._signed_asset(gs_uri, f"page-{page_number}-img-{image_index}"))
                    if len(images) >= max_images:
                        break
                if len(images) >= max_images:
                    break

        elif suffix in {".mp4", ".mov", ".webm", ".mkv"}:
            images.extend(self._extract_video_frames(source_id_norm, source_name, content, max_images))

        return {"source_id": doc["source_id"], "images": images}

    def _extract_video_frames(
        self,
        source_id: str,
        source_name: str,
        content: bytes,
        max_images: int,
    ) -> list[dict]:
        ffmpeg_path = shutil.which("ffmpeg")
        if not ffmpeg_path:
            self.logger.warning("Video frame extraction skipped: ffmpeg not found in PATH.")
            return []

        suffix = Path(source_name).suffix.lower() or ".mp4"
        with tempfile.TemporaryDirectory(prefix="video-frames-") as tmpdir:
            input_path = Path(tmpdir) / f"input{suffix}"
            output_pattern = str(Path(tmpdir) / "frame-%04d.jpg")
            input_path.write_bytes(content)

            # Extract one frame every ~5 seconds, capped by max_images.
            cmd = [
                ffmpeg_path,
                "-hide_banner",
                "-loglevel",
                "error",
                "-i",
                str(input_path),
                "-vf",
                "fps=1/5",
                "-frames:v",
                str(max_images),
                output_pattern,
            ]
            try:
                subprocess.run(cmd, check=True, capture_output=True)
            except subprocess.CalledProcessError as err:
                self.logger.warning(
                    "Video frame extraction failed source_id=%s error=%s",
                    source_id,
                    err.stderr.decode("utf-8", errors="ignore") if err.stderr else str(err),
                )
                return []

            frame_paths = sorted(Path(tmpdir).glob("frame-*.jpg"))
            images: list[dict] = []
            for idx, frame_path in enumerate(frame_paths[:max_images], start=1):
                blob_path = f"extracted/{source_id}/video-frame-{idx:04d}.jpg"
                gs_uri = self._upload_extracted(blob_path, frame_path.read_bytes())
                images.append(self._signed_asset(gs_uri, f"video-frame-{idx:04d}"))
            return images

    def _embed_documents(self, texts: list[str]) -> list[list[float]]:
        vectors: list[list[float]] = []
        batch_size = max(1, EMBED_BATCH_SIZE)
        for start in range(0, len(texts), batch_size):
            batch = texts[start : start + batch_size]
            inputs = [TextEmbeddingInput(text=t, task_type="RETRIEVAL_DOCUMENT") for t in batch]
            outputs = self._run_with_ai_retries(
                lambda: self.embedder.get_embeddings(inputs),
                operation=f"embedding documents batch={len(batch)}",
            )
            vectors.extend([o.values for o in outputs])
        return vectors

    def _embed_query(self, text: str) -> list[float]:
        inputs = [TextEmbeddingInput(text=text, task_type="RETRIEVAL_QUERY")]
        output = self._run_with_ai_retries(
            lambda: self.embedder.get_embeddings(inputs)[0],
            operation="embedding query",
        )
        return output.values

    def _drive_service(self):
        creds, _ = google_auth_default(
            scopes=["https://www.googleapis.com/auth/drive.readonly"]
        )
        return build("drive", "v3", credentials=creds, cache_discovery=False)

    def _drive_file_meta(self, drive, file_id: str) -> dict:
        return (
            drive.files()
            .get(
                fileId=file_id,
                fields="id,name,mimeType",
                supportsAllDrives=True,
            )
            .execute()
        )

    def _drive_list_folder(self, drive, folder_id: str, max_files: int) -> list[dict]:
        query = f"'{folder_id}' in parents and trashed=false"
        items = []
        page_token = None
        while len(items) < max_files:
            resp = (
                drive.files()
                .list(
                    q=query,
                    fields="nextPageToken,files(id,name,mimeType)",
                    includeItemsFromAllDrives=True,
                    supportsAllDrives=True,
                    pageSize=min(100, max_files - len(items)),
                    pageToken=page_token,
                )
                .execute()
            )
            items.extend(resp.get("files", []))
            page_token = resp.get("nextPageToken")
            if not page_token:
                break
        return items

    def _normalize_drive_file(self, item: dict) -> str | None:
        mime = item["mimeType"]
        name = item["name"]
        if mime == "application/vnd.google-apps.presentation":
            return f"{name}.pptx" if not name.lower().endswith(".pptx") else name
        if mime == "application/pdf":
            return name if name.lower().endswith(".pdf") else f"{name}.pdf"
        if (
            mime
            == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        ):
            return name if name.lower().endswith(".pptx") else f"{name}.pptx"
        media_suffix = {
            "image/png": ".png",
            "image/jpeg": ".jpg",
            "image/webp": ".webp",
            "audio/mpeg": ".mp3",
            "audio/wav": ".wav",
            "audio/x-wav": ".wav",
            "audio/mp4": ".m4a",
            "video/mp4": ".mp4",
            "video/quicktime": ".mov",
            "video/webm": ".webm",
            "video/x-matroska": ".mkv",
        }.get(mime)
        if media_suffix:
            return name if name.lower().endswith(media_suffix) else f"{name}{media_suffix}"
        return None

    def _download_drive_file(self, drive, item: dict, target_name: str) -> tuple[str, bytes]:
        mime = item["mimeType"]
        file_id = item["id"]
        if mime == "application/vnd.google-apps.presentation":
            req = (
                drive.files()
                .export_media(
                    fileId=file_id,
                    mimeType="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )
            )
        else:
            req = drive.files().get_media(fileId=file_id, supportsAllDrives=True)

        buffer = io.BytesIO()
        downloader = MediaIoBaseDownload(buffer, req)
        done = False
        while not done:
            _, done = downloader.next_chunk()
        return target_name, buffer.getvalue()

    def _transcribe_media_with_timestamps(
        self,
        data: bytes,
        mime_type: str,
        include_visual_analysis: bool,
        source_uri: str | None = None,
    ) -> dict:
        visual_req = (
            "Also include key visual context affecting speech style or meaning."
            if include_visual_analysis
            else "Ignore visual details and focus on speech only."
        )
        prompt = textwrap.dedent(
            f"""
            Analyze this media and return valid JSON only.
            {visual_req}
            JSON schema:
            {{
              "summary": "short summary",
              "segments": [
                {{
                  "start_seconds": 0.0,
                  "end_seconds": 2.5,
                  "text": "spoken transcript for this segment",
                  "style": "tone and delivery, e.g. excited/neutral/hesitant/whispering",
                  "words": [
                    {{
                      "word": "hello",
                      "start_seconds": 0.1,
                      "end_seconds": 0.5,
                      "style": "tone for this word"
                    }}
                  ]
                }}
              ]
            }}
            Requirements:
            - Keep timestamps monotonic and in seconds.
            - Split into practical segments (~3-12 seconds when possible).
            - Include as many word timings as possible.
            """
        ).strip()
        try:
            media_part = self._media_part(data, mime_type, source_uri=source_uri, prefer_uri=True)
            response = self._generate_with_model_fallback(
                [prompt, media_part],
                operation="media transcription",
                use_media_models=True,
            )
            raw = response.text if hasattr(response, "text") and response.text else str(response)
            parsed = self._parse_json_response(raw)
            if isinstance(parsed, dict):
                return parsed
        except Exception:
            pass

        text = self._describe_binary(
            data,
            mime_type,
            "Transcribe this media and include rough timestamp ranges when possible.",
            source_uri=source_uri,
            prefer_uri=True,
        )
        return {
            "summary": "",
            "segments": [{"start_seconds": 0.0, "end_seconds": None, "text": text, "style": None, "words": []}],
            "text": text,
        }

    def _parse_json_response(self, text: str) -> dict | list:
        clean = text.strip()
        fenced = re.search(r"```(?:json)?\s*(.*?)```", clean, flags=re.DOTALL | re.IGNORECASE)
        if fenced:
            clean = fenced.group(1).strip()
        try:
            return json.loads(clean)
        except json.JSONDecodeError:
            start = clean.find("{")
            end = clean.rfind("}")
            if start >= 0 and end > start:
                return json.loads(clean[start : end + 1])
            raise

    def _describe_binary(
        self,
        data: bytes,
        mime_type: str,
        prompt: str,
        source_uri: str | None = None,
        prefer_uri: bool = False,
    ) -> str:
        media_part = self._media_part(data, mime_type, source_uri=source_uri, prefer_uri=prefer_uri)
        response = self._generate_with_model_fallback(
            [prompt, media_part],
            operation="content generation",
            use_media_models=False,
        )
        text = response.text if hasattr(response, "text") and response.text else str(response)
        return text.strip()

    def _generate_with_model_fallback(self, content, operation: str, use_media_models: bool):
        model_chain = self.media_models if use_media_models else self.text_models
        last_error: Exception | None = None
        for idx, (model_name, model) in enumerate(model_chain, start=1):
            try:
                self.logger.info(
                    "Vertex model selected operation=%s model=%s chain_index=%d/%d",
                    operation,
                    model_name,
                    idx,
                    len(model_chain),
                )
                return self._run_with_ai_retries(
                    lambda: model.generate_content(content),
                    operation=f"{operation} model={model_name}",
                )
            except Exception as err:
                last_error = err
                if idx < len(model_chain) and (
                    self._is_retryable_ai_error(err) or self._is_model_unavailable_error(err)
                ):
                    self.logger.warning(
                        "Switching fallback model operation=%s from=%s to=%s",
                        operation,
                        model_name,
                        model_chain[idx][0],
                    )
                    continue
                raise
        if last_error:
            raise last_error
        raise RuntimeError(f"No model available for operation={operation}")

    def _run_with_ai_retries(self, fn, operation: str):
        delay = max(0.1, VERTEX_RETRY_INITIAL_SECONDS)
        last_error: Exception | None = None
        total_attempts = max(1, VERTEX_MAX_RETRIES + 1)
        for attempt in range(1, total_attempts + 1):
            self.logger.info("Vertex call start operation=%s attempt=%d/%d", operation, attempt, total_attempts)
            try:
                result = fn()
                self.logger.info("Vertex call success operation=%s attempt=%d/%d", operation, attempt, total_attempts)
                return result
            except Exception as err:
                last_error = err
                self.logger.warning(
                    "Vertex call failed operation=%s attempt=%d/%d error=%s",
                    operation,
                    attempt,
                    total_attempts,
                    err,
                )
                if not self._is_retryable_ai_error(err) or attempt == total_attempts:
                    break
                sleep_for = min(delay, VERTEX_RETRY_MAX_SECONDS)
                sleep_for += random.uniform(0.0, min(1.0, sleep_for * 0.2))
                self.logger.info(
                    "Vertex retry scheduled operation=%s next_attempt=%d/%d sleep_seconds=%.2f",
                    operation,
                    attempt + 1,
                    total_attempts,
                    sleep_for,
                )
                time.sleep(sleep_for)
                delay = min(delay * 2, VERTEX_RETRY_MAX_SECONDS)

        if last_error and self._is_retryable_ai_error(last_error):
            raise RuntimeError(
                f"Vertex AI rate limit or temporary capacity issue during {operation}. "
                "Retries exhausted; please try again in a few minutes."
            ) from last_error
        if last_error:
            raise last_error
        raise RuntimeError(f"Unknown AI error during {operation}.")

    def _is_retryable_ai_error(self, err: Exception) -> bool:
        if isinstance(
            err,
            (
                google_api_exceptions.ResourceExhausted,
                google_api_exceptions.TooManyRequests,
                google_api_exceptions.ServiceUnavailable,
                google_api_exceptions.DeadlineExceeded,
            ),
        ):
            return True
        msg = str(err).lower()
        return (
            "429" in msg
            or "resource exhausted" in msg
            or "quota" in msg
            or "rate limit" in msg
            or "temporarily unavailable" in msg
        )

    def _is_model_unavailable_error(self, err: Exception) -> bool:
        if isinstance(
            err,
            (
                google_api_exceptions.NotFound,
                google_api_exceptions.PermissionDenied,
            ),
        ):
            return True
        msg = str(err).lower()
        return (
            "publisher model" in msg
            and ("not found" in msg or "does not have access" in msg)
        )

    def _build_model_chain(self, primary_model: str, fallback_models: list[str]) -> list[tuple[str, GenerativeModel]]:
        seen: set[str] = set()
        chain: list[tuple[str, GenerativeModel]] = []
        for model_name in [primary_model, *fallback_models]:
            if not model_name or model_name in seen:
                continue
            seen.add(model_name)
            chain.append((model_name, GenerativeModel(model_name)))
        return chain

    def _media_part(
        self,
        data: bytes,
        mime_type: str,
        source_uri: str | None = None,
        prefer_uri: bool = False,
    ) -> Part:
        if prefer_uri and source_uri and source_uri.startswith("gs://"):
            return Part.from_uri(source_uri, mime_type)
        return Part.from_data(data=data, mime_type=mime_type)

    def _maybe_sign_gs_uri(self, gs_uri: str) -> str | None:
        if not gs_uri.startswith("gs://"):
            return None
        try:
            return self._signed_asset(gs_uri, "source").get("signed_url")
        except Exception:
            return None

    def _format_timecode(self, seconds: float | None) -> str | None:
        if seconds is None:
            return None
        total = max(0.0, float(seconds))
        hrs = int(total // 3600)
        mins = int((total % 3600) // 60)
        secs = total % 60
        return f"{hrs:02d}:{mins:02d}:{secs:06.3f}"

    def _build_playback_url(self, base_url: str | None, start_seconds: float | None) -> str | None:
        if not base_url:
            return None
        if start_seconds is None:
            return base_url
        return f"{base_url}#t={max(0.0, float(start_seconds)):.3f}"

    def _safe_float(self, value) -> float | None:
        if value is None or value == "":
            return None
        try:
            return float(value)
        except (TypeError, ValueError):
            return None

    def _extract_date_from_text(self, text: str) -> str | None:
        if not text:
            return None
        for pattern in [
            r"\b(20\d{2})-(\d{2})-(\d{2})\b",
            r"\b(20\d{2})/(\d{2})/(\d{2})\b",
            r"\b(\d{1,2})/(\d{1,2})/(20\d{2})\b",
        ]:
            m = re.search(pattern, text)
            if not m:
                continue
            try:
                if pattern.startswith(r"\b(\d{1,2})"):
                    mm, dd, yyyy = int(m.group(1)), int(m.group(2)), int(m.group(3))
                else:
                    yyyy, mm, dd = int(m.group(1)), int(m.group(2)), int(m.group(3))
                return date(yyyy, mm, dd).isoformat()
            except ValueError:
                pass

        month_match = re.search(
            r"\b(january|february|march|april|may|june|july|august|september|october|november|december)\s+(20\d{2})\b",
            text,
            flags=re.IGNORECASE,
        )
        if month_match:
            month = {
                "january": 1,
                "february": 2,
                "march": 3,
                "april": 4,
                "may": 5,
                "june": 6,
                "july": 7,
                "august": 8,
                "september": 9,
                "october": 10,
                "november": 11,
                "december": 12,
            }[month_match.group(1).lower()]
            year = int(month_match.group(2))
            return date(year, month, 1).isoformat()
        return None

    def _extract_date_filter(self, question: str) -> tuple[str | None, str | None]:
        q = question.lower()
        year_match = re.search(r"\b(20\d{2})\b", q)
        year = int(year_match.group(1)) if year_match else None

        season_match = re.search(r"\b(beginning of\s+)?(spring|summer|fall|autumn|winter)\s+(20\d{2})\b", q)
        if season_match:
            beginning = bool(season_match.group(1))
            season = season_match.group(2)
            yr = int(season_match.group(3))
            ranges = {
                "spring": (date(yr, 3, 1), date(yr, 5, 31)),
                "summer": (date(yr, 6, 1), date(yr, 8, 31)),
                "fall": (date(yr, 9, 1), date(yr, 11, 30)),
                "autumn": (date(yr, 9, 1), date(yr, 11, 30)),
                "winter": (date(yr, 12, 1), date(yr + 1, 2, 28)),
            }
            start, end = ranges[season]
            if beginning:
                end = min(
                    end,
                    start.replace(day=1) + (date(start.year, start.month, 28) - date(start.year, start.month, 1)),
                )
                end = date(start.year, min(start.month + 1, 12), 15)
            return start.isoformat(), end.isoformat()

        month_year = re.search(
            r"\b(january|february|march|april|may|june|july|august|september|october|november|december)\s+(20\d{2})\b",
            q,
        )
        if month_year:
            month = {
                "january": 1,
                "february": 2,
                "march": 3,
                "april": 4,
                "may": 5,
                "june": 6,
                "july": 7,
                "august": 8,
                "september": 9,
                "october": 10,
                "november": 11,
                "december": 12,
            }[month_year.group(1)]
            yr = int(month_year.group(2))
            start = date(yr, month, 1)
            end = date(yr, month, 28)
            return start.isoformat(), end.isoformat()

        if year and ("in " + str(year) in q or "from " + str(year) in q):
            return date(year, 1, 1).isoformat(), date(year, 12, 31).isoformat()
        return None, None

    def _source_id_from_uri(self, source_uri: str) -> str:
        return uuid.uuid5(uuid.NAMESPACE_URL, source_uri).hex

    def _download_gcs_uri(self, gs_uri: str) -> bytes:
        bucket_name, blob_name = self._parse_gs_uri(gs_uri)
        blob = self.storage.bucket(bucket_name).blob(blob_name)
        return blob.download_as_bytes()

    def _upload_extracted(self, blob_path: str, data: bytes) -> str:
        bucket = self.storage.bucket(GCS_BUCKET)
        blob = bucket.blob(blob_path)
        if not blob.exists():
            blob.upload_from_string(data)
        return f"gs://{GCS_BUCKET}/{blob_path}"

    def _signed_asset(self, gs_uri: str, label: str) -> dict:
        bucket_name, blob_name = self._parse_gs_uri(gs_uri)
        blob = self.storage.bucket(bucket_name).blob(blob_name)
        url = blob.generate_signed_url(version="v4", expiration=timedelta(hours=1), method="GET")
        return {"label": label, "gs_uri": gs_uri, "signed_url": url}

    def _parse_gs_uri(self, gs_uri: str) -> tuple[str, str]:
        parsed = urlparse(gs_uri)
        if parsed.scheme != "gs":
            raise ValueError(f"Invalid gs URI: {gs_uri}")
        return parsed.netloc, parsed.path.lstrip("/")
